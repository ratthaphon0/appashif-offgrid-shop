"""Build the teaching deck for AppAshif: Product data from GitHub JSON."""

import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "deliverables" / "AppAshif_Product_GitHub_JSON_Lesson.pptx"
JSON_PATH = ROOT / "assets" / "data" / "products.json"

W, H = Inches(13.333), Inches(7.5)
INK = "111111"
CREAM = "F6F0E4"
PAPER = "FFFCF5"
NEON = "C8FF35"
PINK = "FF4FA3"
PURPLE = "8B5CF6"
ORANGE = "FF6B2C"
MUTED = "6E685F"


def rgb(value):
    return RGBColor.from_string(value.lstrip("#"))


def shape(slide, x, y, w, h, fill, line=INK, radius=True):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    item = slide.shapes.add_shape(kind, x, y, w, h)
    item.fill.solid()
    item.fill.fore_color.rgb = rgb(fill)
    item.line.color.rgb = rgb(line)
    item.line.width = Pt(1.5)
    return item


def text(slide, value, x, y, w, h, size=18, color=INK, bold=False, font="Aptos", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = value
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def code(slide, value, x, y, w, h, size=11):
    card = shape(slide, x, y, w, h, "202124", "202124", False)
    card.line.width = Pt(0)
    box = text(slide, value, x + Inches(.18), y + Inches(.16), w - Inches(.36), h - Inches(.28), size, "E8EAED", False, "Courier New")
    box.text_frame.paragraphs[0].space_after = Pt(0)
    return card


def header(slide, number, title, kicker="INTERNET PROGRAMMING · CONTINUED"):
    text(slide, kicker, Inches(.55), Inches(.34), Inches(7), Inches(.25), 10, PURPLE, True)
    text(slide, title, Inches(.55), Inches(.64), Inches(11.6), Inches(.62), 28, INK, True)
    tag = shape(slide, Inches(12.05), Inches(.36), Inches(.7), Inches(.42), NEON)
    text(slide, f"{number:02}", Inches(12.05), Inches(.45), Inches(.7), Inches(.2), 10, INK, True, align=PP_ALIGN.CENTER)
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(.55), Inches(1.36), Inches(12.2), Inches(.03)).fill.solid()
    line = slide.shapes[-1]
    line.fill.fore_color.rgb = rgb(INK)
    line.line.fill.background()


def footer(slide):
    text(slide, "AppAshif OFF//GRID · Ratthaphon Khan · SEC 800", Inches(.55), Inches(7.12), Inches(7), Inches(.18), 8, MUTED, True)
    text(slide, "Expo 57 · React Native · GitHub JSON", Inches(9.1), Inches(7.12), Inches(3.65), Inches(.18), 8, MUTED, True, align=PP_ALIGN.RIGHT)


def bullet(slide, value, x, y, w, color=INK):
    text(slide, "•", x, y, Inches(.18), Inches(.25), 16, PURPLE, True)
    text(slide, value, x + Inches(.23), y, w - Inches(.23), Inches(.45), 14, color, False)


def add_title(slide):
    shape(slide, Inches(.55), Inches(.48), Inches(12.2), Inches(6.5), NEON, NEON, False).line.width = Pt(0)
    shape(slide, Inches(7.92), Inches(.48), Inches(4.83), Inches(6.5), PURPLE, PURPLE, False).line.width = Pt(0)
    text(slide, "INTERNET PROGRAMMING", Inches(.85), Inches(.9), Inches(5.8), Inches(.3), 12, INK, True)
    text(slide, "CONTINUE FROM\nLAST CLASS", Inches(.85), Inches(1.45), Inches(6.5), Inches(1.5), 39, INK, True)
    text(slide, "Product catalogue from GitHub JSON\n→ validate → render → recover safely", Inches(.87), Inches(3.22), Inches(5.9), Inches(.7), 18, INK, True)
    shape(slide, Inches(.85), Inches(4.33), Inches(4.1), Inches(.55), INK, INK)
    text(slide, "EXPO 57 / REACT NATIVE", Inches(.98), Inches(4.49), Inches(3.82), Inches(.18), 11, NEON, True)
    text(slide, "Today’s goal: make the existing OFF//GRID product page use data that has been committed to GitHub.", Inches(.87), Inches(5.32), Inches(5.7), Inches(.62), 15, INK, False)
    text(slide, "GITHUB\nJSON", Inches(8.45), Inches(1.35), Inches(3.4), Inches(1.2), 42, "FFFFFF", True, align=PP_ALIGN.CENTER)
    shape(slide, Inches(8.7), Inches(3.28), Inches(3.25), Inches(2.02), PAPER, INK)
    text(slide, "products.json\n4 detailed products\nvalidated before UI", Inches(8.92), Inches(3.68), Inches(2.82), Inches(.9), 18, INK, True, align=PP_ALIGN.CENTER)
    text(slide, "Ratthaphon Khan · 6730202734 · SEC 800", Inches(.87), Inches(6.43), Inches(5.8), Inches(.2), 10, INK, True)


def build():
    products = json.loads(JSON_PATH.read_text())
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb(CREAM)
    add_title(slide)

    slide = prs.slides.add_slide(blank); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb(CREAM)
    header(slide, 2, "Where the product data lives")
    text(slide, "Repository contract", Inches(.72), Inches(1.75), Inches(4), Inches(.3), 18, INK, True)
    for i, line in enumerate(["assets/data/products.json", "↓ commit + push to main", "raw.githubusercontent.com/.../products.json", "↓ fetch in the Expo app", "ProductCard renders each item"]):
        y = Inches(2.22 + i * .65)
        fill = NEON if i in (0, 4) else PAPER
        shape(slide, Inches(.82), y, Inches(4.6), Inches(.42), fill)
        text(slide, line, Inches(1.02), y + Inches(.12), Inches(4.15), Inches(.16), 12, INK, True, "Courier New")
    code(slide, "export const PRODUCT_JSON_URL =\n  'https://raw.githubusercontent.com/ratthaphon0/\n   appashif-offgrid-shop/main/assets/data/products.json';", Inches(6.18), Inches(1.76), Inches(5.9), Inches(1.48), 11)
    bullet(slide, "A raw GitHub URL returns JSON directly; it is a simple catalogue API for this class.", Inches(6.23), Inches(3.7), Inches(5.5))
    bullet(slide, "The URL points to the main branch, so students see the latest committed data after Refresh.", Inches(6.23), Inches(4.35), Inches(5.5))
    bullet(slide, "Keep images, price, stock, sale price, and display metadata together in one product record.", Inches(6.23), Inches(5.0), Inches(5.5))
    footer(slide)

    slide = prs.slides.add_slide(blank); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb(CREAM)
    header(slide, 3, "Detailed JSON: one source of truth")
    sample = json.dumps(products[0], indent=2, ensure_ascii=False)
    code(slide, sample, Inches(.7), Inches(1.66), Inches(6.15), Inches(5.12), 9.5)
    text(slide, "Fields used in the UI", Inches(7.25), Inches(1.72), Inches(4), Inches(.3), 18, INK, True)
    for i, item in enumerate([("identity", "id, name, category"), ("commerce", "price, originalPrice, stock"), ("presentation", "badge, colors, edition"), ("media", "images: label + uri"), ("content", "description for detail views")]):
        y = Inches(2.28 + i * .73)
        shape(slide, Inches(7.25), y, Inches(4.92), Inches(.52), PAPER)
        text(slide, item[0].upper(), Inches(7.47), y + Inches(.1), Inches(1.32), Inches(.18), 10, PURPLE, True)
        text(slide, item[1], Inches(8.82), y + Inches(.1), Inches(3.05), Inches(.22), 12, INK, True)
    text(slide, "Rule: update the JSON once, then every screen receives the same product information.", Inches(7.25), Inches(6.24), Inches(4.92), Inches(.42), 13, INK, True)
    footer(slide)

    slide = prs.slides.add_slide(blank); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb(CREAM)
    header(slide, 4, "Fetch, validate, then render")
    code(slide, "const response = await fetch(PRODUCT_JSON_URL);\nif (!response.ok) throw new Error(`GitHub returned ${response.status}`);\n\nconst data: unknown = await response.json();\nif (!Array.isArray(data) || !data.every(isProduct)) {\n  throw new Error('GitHub JSON does not match the product schema');\n}\n\nsetProducts(data);\nsetSource('github');", Inches(.75), Inches(1.75), Inches(6.28), Inches(4.62), 11)
    shape(slide, Inches(7.62), Inches(1.9), Inches(4.48), Inches(.72), NEON)
    text(slide, "1. FETCH", Inches(7.86), Inches(2.14), Inches(4), Inches(.2), 15, INK, True)
    shape(slide, Inches(7.62), Inches(2.95), Inches(4.48), Inches(.72), "E9DCFF")
    text(slide, "2. VALIDATE SHAPE", Inches(7.86), Inches(3.19), Inches(4), Inches(.2), 15, INK, True)
    shape(slide, Inches(7.62), Inches(4.0), Inches(4.48), Inches(.72), PINK)
    text(slide, "3. UPDATE CONTEXT", Inches(7.86), Inches(4.24), Inches(4), Inches(.2), 15, INK, True)
    shape(slide, Inches(7.62), Inches(5.05), Inches(4.48), Inches(.72), PAPER)
    text(slide, "4. PRODUCTCARD MAPS DATA", Inches(7.86), Inches(5.29), Inches(4), Inches(.2), 15, INK, True)
    text(slide, "`ProductProvider` keeps fetching logic out of screen components.", Inches(7.62), Inches(6.18), Inches(4.5), Inches(.32), 12, MUTED, True)
    footer(slide)

    slide = prs.slides.add_slide(blank); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb(CREAM)
    header(slide, 5, "Use the shared catalogue in the Products page")
    code(slide, "const { products, source, isLoading, error, refreshProducts } = useProducts();\n\n<PageHeading\n  title=\"PRODUCTS\"\n  badge={isLoading ? 'SYNCING…' : `${products.length} ITEMS`}\n/>\n\n{products.map((product) => (\n  <ProductCard key={product.id} product={product}\n    width={cardWidth} onAdd={() => addItem(product.id)} />\n))}", Inches(.72), Inches(1.7), Inches(6.1), Inches(4.7), 11)
    shape(slide, Inches(7.28), Inches(1.72), Inches(4.75), Inches(4.72), PAPER)
    text(slide, "PRODUCTS", Inches(7.62), Inches(2.02), Inches(2.7), Inches(.35), 25, INK, True)
    shape(slide, Inches(10.25), Inches(1.94), Inches(1.38), Inches(.38), PINK)
    text(slide, "4 ITEMS", Inches(10.25), Inches(2.07), Inches(1.38), Inches(.15), 9, INK, True, align=PP_ALIGN.CENTER)
    shape(slide, Inches(7.58), Inches(2.62), Inches(4.12), Inches(.42), NEON)
    text(slide, "LIVE: GITHUB JSON                     REFRESH", Inches(7.76), Inches(2.75), Inches(3.75), Inches(.17), 10, INK, True)
    for i, product in enumerate(products[:2]):
        x = Inches(7.58 + i * 2.04)
        shape(slide, x, Inches(3.36), Inches(1.78), Inches(2.37), "FFFFFF")
        shape(slide, x + Inches(.12), Inches(3.53), Inches(1.54), Inches(.83), product["imageColor"], INK)
        text(slide, product["category"].upper(), x + Inches(.14), Inches(4.55), Inches(1.5), Inches(.15), 7, MUTED, True)
        text(slide, product["name"], x + Inches(.14), Inches(4.77), Inches(1.48), Inches(.4), 11, INK, True)
        text(slide, f"฿{product['price']:,} · {product['stock']} left", x + Inches(.14), Inches(5.34), Inches(1.48), Inches(.15), 8, PURPLE, True)
    text(slide, "One hook supplies the same array to Home, Products, and Cart.", Inches(7.5), Inches(6.08), Inches(4.35), Inches(.25), 12, INK, True, align=PP_ALIGN.CENTER)
    footer(slide)

    slide = prs.slides.add_slide(blank); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb(CREAM)
    header(slide, 6, "Make failure safe, not invisible")
    shape(slide, Inches(.78), Inches(1.78), Inches(5.38), Inches(4.68), "FFE7EE")
    text(slide, "If GitHub is unavailable", Inches(1.08), Inches(2.12), Inches(4.6), Inches(.35), 22, INK, True)
    bullet(slide, "Catch network / HTTP errors.", Inches(1.1), Inches(2.8), Inches(4.5))
    bullet(slide, "Reject malformed JSON before it reaches the UI.", Inches(1.1), Inches(3.42), Inches(4.5))
    bullet(slide, "Show the local fallback catalogue.", Inches(1.1), Inches(4.04), Inches(4.5))
    bullet(slide, "Tell the user: OFFLINE: LOCAL FALLBACK.", Inches(1.1), Inches(4.66), Inches(4.5))
    shape(slide, Inches(7.0), Inches(1.78), Inches(5.38), Inches(4.68), "EAFFD4")
    text(slide, "Why this is better", Inches(7.3), Inches(2.12), Inches(4.6), Inches(.35), 22, INK, True)
    bullet(slide, "The product grid never becomes blank.", Inches(7.32), Inches(2.8), Inches(4.45))
    bullet(slide, "A broken JSON edit cannot crash the cart.", Inches(7.32), Inches(3.42), Inches(4.45))
    bullet(slide, "The status bar makes the data source testable in class.", Inches(7.32), Inches(4.04), Inches(4.45))
    text(slide, "Demo test: turn off Wi-Fi → reload → products still render.", Inches(7.32), Inches(5.18), Inches(4.5), Inches(.4), 14, PURPLE, True)
    footer(slide)

    slide = prs.slides.add_slide(blank); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb(CREAM)
    header(slide, 7, "Class demo checklist")
    items = [
        ("1", "Edit products.json", "Change one product name, stock, or price."),
        ("2", "Commit and push", "GitHub becomes the published catalogue."),
        ("3", "Open Products", "Press REFRESH and see LIVE: GITHUB JSON."),
        ("4", "Add to cart", "The selected product comes from the same shared data."),
        ("5", "Test offline", "Fallback protects the UI when GitHub is unreachable."),
    ]
    for i, (number, title, detail) in enumerate(items):
        y = Inches(1.66 + i * .91)
        shape(slide, Inches(.9), y, Inches(.55), Inches(.55), NEON if i % 2 == 0 else PINK)
        text(slide, number, Inches(.9), y + Inches(.18), Inches(.55), Inches(.15), 12, INK, True, align=PP_ALIGN.CENTER)
        text(slide, title, Inches(1.75), y + Inches(.04), Inches(2.4), Inches(.22), 16, INK, True)
        text(slide, detail, Inches(4.2), y + Inches(.07), Inches(6.8), Inches(.22), 13, MUTED, False)
    shape(slide, Inches(9.94), Inches(5.96), Inches(2.1), Inches(.56), PURPLE, PURPLE)
    text(slide, "Q / A", Inches(9.94), Inches(6.15), Inches(2.1), Inches(.18), 15, "FFFFFF", True, align=PP_ALIGN.CENTER)
    footer(slide)

    OUT.parent.mkdir(exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
